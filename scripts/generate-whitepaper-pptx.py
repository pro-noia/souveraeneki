#!/usr/bin/env python3
"""
Generiert public/whitepaper.pptx aus souveraene-ki-whitepaper.docx.

Brand-Light-Theme: weißer Hintergrund (druckfreundlich), dunkler Headline-Text,
Petrol-Akzent (#3a90aa), Manrope für Display + Cousine-Mono für Eyebrows.

Quellen:
- Source-Doc: souveraene-ki-whitepaper.docx (Repo-Root)
- Brand-Tokens: src/app/globals.css (OKLCH dort, hier als sRGB hex)
- Visual-Referenz: src/lib/og-image.tsx

Aufruf:  npm run pptx:whitepaper    (siehe scripts/WHITEPAPER-PPTX.md)
"""

from __future__ import annotations

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------
# Brand
# --------------------------------------------------------------------------

# sRGB-Approximationen der OKLCH-Tokens aus globals.css, adaptiert für hellen
# Hintergrund. pptx unterstützt OKLCH nicht — daher Hex.
BG          = RGBColor(0xFF, 0xFF, 0xFF)  # pures Weiß, druck-clean
HEADLINE    = RGBColor(0x0C, 0x16, 0x26)  # tinte-deep
BODY        = RGBColor(0x1F, 0x29, 0x37)  # anthrazit (lesbar auf weiß)
BODY_MUTED  = RGBColor(0x4B, 0x55, 0x63)  # heller für Captions
PETROL      = RGBColor(0x3A, 0x90, 0xAA)  # petrol-bright — Eyebrows + Akzent
PETROL_DEEP = RGBColor(0x2A, 0x6C, 0x7E)
MONDSTEIN   = RGBColor(0x55, 0x7B, 0x95)  # mondstein-deep — sek. Caption
ROST        = RGBColor(0xC4, 0x5E, 0x3A)  # warning
SALBEI      = RGBColor(0x7F, 0xA8, 0x73)  # success
BERNSTEIN   = RGBColor(0xC4, 0x92, 0x3A)  # caution
LINE_SUBTLE = RGBColor(0xE2, 0xDF, 0xD8)  # sehr leichte Trennlinie (Pergament-Off)

# Fonts — Manrope/Cousine sind die Site-Brand. Falls beim Empfänger nicht
# installiert: PowerPoint/Keynote fallen auf Helvetica Neue / Menlo zurück.
FONT_DISPLAY = "Manrope"
FONT_MONO    = "Cousine"

# Slide-Dimensions (16:9 standard für pptx ist 13.333" × 7.5")
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Margins
M_TOP    = Inches(0.65)
M_BOTTOM = Inches(0.5)
M_LEFT   = Inches(0.7)
M_RIGHT  = Inches(0.7)

# --------------------------------------------------------------------------
# Slide-Master-Helpers
# --------------------------------------------------------------------------


def add_blank_slide(prs: Presentation):
    """Eine leere Slide mit Brand-Master-Elementen (Petrol-Bar oben, Logo, Footer)."""
    blank_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(blank_layout)

    # Hintergrund: pures Weiß (überschreibt eventuelle Layout-Defaults).
    bg_rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    bg_rect.line.fill.background()
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = BG
    bg_rect.shadow.inherit = False

    # Petrol-Top-Bar (3 pt) — echo zum OG-Image.
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.04)
    )
    top_bar.line.fill.background()
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PETROL

    # Logo oben links: Petrol-Square + Wordmark.
    logo_y = Inches(0.18)
    sq = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), logo_y, Inches(0.18), Inches(0.18)
    )
    sq.line.fill.background()
    sq.fill.solid()
    sq.fill.fore_color.rgb = PETROL

    wm = slide.shapes.add_textbox(
        Inches(0.95), Inches(0.12), Inches(3), Inches(0.3)
    )
    set_text(
        wm,
        "Souveräne KI",
        font=FONT_DISPLAY,
        size=11,
        bold=True,
        color=HEADLINE,
        spacing=0.01,
    )

    return slide


def add_footer(slide, page_num: int, total_pages: int):
    """Seitenzahl + Brand-Footer unten."""
    # Linke Footer-Linie: Stand
    left = slide.shapes.add_textbox(
        Inches(0.7), Inches(7.1), Inches(6), Inches(0.3)
    )
    set_text(
        left,
        "PRAXISLEITFADEN · SOUVERÄNE KI · MAI 2026",
        font=FONT_MONO,
        size=7,
        color=MONDSTEIN,
        spacing=0.18,
    )

    # Rechte Seitenzahl
    right = slide.shapes.add_textbox(
        Inches(11), Inches(7.1), Inches(1.633), Inches(0.3)
    )
    set_text(
        right,
        f"{page_num:02d} · {total_pages:02d}",
        font=FONT_MONO,
        size=7,
        color=MONDSTEIN,
        align=PP_ALIGN.RIGHT,
        spacing=0.12,
    )


# --------------------------------------------------------------------------
# Text-Helpers
# --------------------------------------------------------------------------


def set_text(
    shape,
    text: str,
    *,
    font: str = FONT_DISPLAY,
    size: float = 14,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = BODY,
    align=PP_ALIGN.LEFT,
    spacing: float = 0.0,
    line_spacing: float | None = None,
):
    """Setzt Text mit Brand-Formatierung in einer Shape."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    # Letter-spacing in pptx über spc (XML-Attribut) — python-pptx-API hat
    # das nicht direkt; wir setzen es via XML wenn nötig.
    if spacing > 0:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(spacing * 100)))  # in 1/100 pt
    return tf


def add_text(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    **kwargs,
):
    """Convenience: TextBox + set_text in einem."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    set_text(tb, text, **kwargs)
    return tb


def add_paragraph(tf, text: str, **kwargs):
    """Fügt einen weiteren Paragraph zum bestehenden Textframe hinzu."""
    p = tf.add_paragraph()
    p.alignment = kwargs.pop("align", PP_ALIGN.LEFT)
    line_spacing = kwargs.pop("line_spacing", None)
    if line_spacing is not None:
        p.line_spacing = line_spacing
    space_before = kwargs.pop("space_before", None)
    if space_before is not None:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = kwargs.get("font", FONT_DISPLAY)
    f.size = Pt(kwargs.get("size", 14))
    f.bold = kwargs.get("bold", False)
    f.italic = kwargs.get("italic", False)
    f.color.rgb = kwargs.get("color", BODY)
    spacing = kwargs.get("spacing", 0.0)
    if spacing > 0:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(spacing * 100)))
    return p


def hr_line(slide, top, left=Inches(0.7), width=Inches(11.9), color=LINE_SUBTLE):
    """Schmale horizontale Trennlinie."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.012)
    )
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = color


def vert_line(slide, left, top, height, color=LINE_SUBTLE):
    """Schmale vertikale Trennlinie."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Inches(0.008), height
    )
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = color


# --------------------------------------------------------------------------
# Slide-Builder
# --------------------------------------------------------------------------


def build_cover(prs):
    slide = add_blank_slide(prs)

    # Eyebrow zentriert oben
    add_text(
        slide,
        "PRAXISLEITFADEN  ·  MAI 2026",
        Inches(0.7),
        Inches(1.8),
        Inches(11.9),
        Inches(0.4),
        font=FONT_MONO,
        size=11,
        color=PETROL,
        spacing=0.30,
        align=PP_ALIGN.CENTER,
    )

    # Riesige Headline
    add_text(
        slide,
        "Souveräne KI in Europa.",
        Inches(0.7),
        Inches(2.4),
        Inches(11.9),
        Inches(1.5),
        font=FONT_DISPLAY,
        size=64,
        bold=True,
        color=HEADLINE,
        spacing=-0.6,
        align=PP_ALIGN.CENTER,
    )

    # Subhead
    add_text(
        slide,
        "Vom Intelligence Taker zum Intelligence Maker —\n"
        "ein Praxisleitfaden für CEOs und CIOs.",
        Inches(2),
        Inches(4.2),
        Inches(9.3),
        Inches(1),
        font=FONT_DISPLAY,
        size=20,
        color=BODY,
        line_spacing=1.35,
        align=PP_ALIGN.CENTER,
    )

    # Drei Akte als Mini-Pills unten
    pills_y = Inches(5.6)
    pill_w = Inches(3.6)
    gap = Inches(0.25)
    total_w = pill_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) / 2

    for i, (num, title, desc) in enumerate(
        [
            ("01", "DER WECKRUF", "Warum Europäische Unternehmen ihre KI heute mehr mieten als besitzen."),
            ("02", "DIE LANDKARTE", "Souveränität, Compliance-Welle, Säulen, Ökosystem."),
            ("03", "DER SPIELZUG", "CSD-Framework, 90-Tage-Roadmap und 12-Fragen-Check."),
        ]
    ):
        x = start_x + i * (pill_w + gap)

        # Petrol-Top-Line der Karte
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, pills_y, pill_w, Inches(0.025)
        )
        line.line.fill.background()
        line.fill.solid()
        line.fill.fore_color.rgb = PETROL

        # Num
        add_text(
            slide,
            num,
            x,
            pills_y + Inches(0.1),
            pill_w,
            Inches(0.35),
            font=FONT_MONO,
            size=13,
            color=PETROL,
            spacing=0.20,
        )
        # Title
        add_text(
            slide,
            title,
            x,
            pills_y + Inches(0.45),
            pill_w,
            Inches(0.4),
            font=FONT_DISPLAY,
            size=15,
            bold=True,
            color=HEADLINE,
            spacing=0.10,
        )
        # Desc
        add_text(
            slide,
            desc,
            x,
            pills_y + Inches(0.85),
            pill_w,
            Inches(0.8),
            font=FONT_DISPLAY,
            size=10,
            color=BODY_MUTED,
            line_spacing=1.35,
        )

    # Footer-Meta
    add_text(
        slide,
        "14 Seiten · 28 Quellen mit URLs · Mai 2026",
        Inches(0.7),
        Inches(7),
        Inches(11.9),
        Inches(0.3),
        font=FONT_MONO,
        size=9,
        color=MONDSTEIN,
        spacing=0.20,
        align=PP_ALIGN.CENTER,
    )

    return slide


def build_thesis(prs):
    """Die Kernthese (vor Executive Summary)."""
    slide = add_blank_slide(prs)

    # Eyebrow
    add_text(
        slide,
        "DIE KERNTHESE",
        Inches(0.7),
        Inches(1.4),
        Inches(11.9),
        Inches(0.4),
        font=FONT_MONO,
        size=12,
        color=PETROL,
        spacing=0.30,
        align=PP_ALIGN.CENTER,
    )

    # Großer Quote-Style-Satz
    tb = slide.shapes.add_textbox(
        Inches(1.3), Inches(2.4), Inches(10.7), Inches(2.5)
    )
    set_text(
        tb,
        "Wer seine KI mietet,",
        font=FONT_DISPLAY,
        size=48,
        bold=True,
        color=HEADLINE,
        align=PP_ALIGN.CENTER,
        spacing=-0.5,
        line_spacing=1.1,
    )
    add_paragraph(
        tb.text_frame,
        "baut auf gemietetem Grund.",
        font=FONT_DISPLAY,
        size=48,
        bold=True,
        color=PETROL,
        align=PP_ALIGN.CENTER,
        line_spacing=1.1,
    )

    # Erläuterung darunter
    add_text(
        slide,
        "Souveränität bedeutet nicht Autarkie — sondern Agency. Drei Schritte. "
        "Vier Branchen. Sieben Aktionen.\n"
        "Ein Praxisleitfaden, der die Souveränitätsfrage operationalisiert.",
        Inches(2),
        Inches(5.4),
        Inches(9.3),
        Inches(1.2),
        font=FONT_DISPLAY,
        size=16,
        color=BODY,
        line_spacing=1.5,
        align=PP_ALIGN.CENTER,
    )

    return slide


def build_executive_summary(prs):
    """Executive Summary mit narrativem Block links + 3 Stats rechts."""
    slide = add_blank_slide(prs)

    # Eyebrow
    add_text(
        slide,
        "EXECUTIVE SUMMARY",
        Inches(0.7),
        Inches(0.9),
        Inches(11.9),
        Inches(0.4),
        font=FONT_MONO,
        size=11,
        color=PETROL,
        spacing=0.30,
    )

    # H1
    add_text(
        slide,
        "Wer seine KI mietet, baut auf gemietetem Grund.",
        Inches(0.7),
        Inches(1.3),
        Inches(11.9),
        Inches(1),
        font=FONT_DISPLAY,
        size=32,
        bold=True,
        color=HEADLINE,
        spacing=-0.3,
        line_spacing=1.15,
    )

    hr_line(slide, Inches(2.55))

    # Linke Spalte: narrativer Auftakt
    left_x = Inches(0.7)
    col_w = Inches(7.5)

    body_tb = slide.shapes.add_textbox(left_x, Inches(2.85), col_w, Inches(3.8))
    set_text(
        body_tb,
        "Es ist Montagmorgen, 6:14 Uhr. In einer Produktionshalle in Schwaben "
        "steht der Geschäftsführer eines Automobilzulieferers vor stillstehenden "
        "Maschinen. Nicht wegen eines technischen Defekts — sondern weil ein "
        "Cloud-Anbieter in San Francisco über Nacht die Terms of Service "
        "geändert hat. Pro Stunde Stillstand: 180.000 Euro.",
        font=FONT_DISPLAY,
        size=12,
        color=BODY,
        line_spacing=1.5,
    )
    add_paragraph(
        body_tb.text_frame,
        "Die Szene ist konstruiert. Jedes ihrer Elemente ist real.",
        font=FONT_DISPLAY,
        size=12,
        italic=True,
        color=BODY,
        line_spacing=1.5,
        space_before=8,
    )
    add_paragraph(
        body_tb.text_frame,
        "Souveränität bedeutet nicht Autarkie. Sie bedeutet Agency — die "
        "Freiheit, Abhängigkeiten bewusst zu wählen, mitzugestalten und "
        "jederzeit zu verlassen.",
        font=FONT_DISPLAY,
        size=12,
        color=BODY,
        line_spacing=1.5,
        space_before=8,
    )
    add_paragraph(
        body_tb.text_frame,
        "▸  Akt I — Der Weckruf. Vier regulierte Branchen, konkrete Schmerzpunkte. "
        "Microsoft France am 10. Juni 2025 unter Eid: bestätigt, was viele ahnten.",
        font=FONT_DISPLAY,
        size=11,
        color=BODY,
        line_spacing=1.4,
        space_before=10,
    )
    add_paragraph(
        body_tb.text_frame,
        "▸  Akt II — Die Landkarte. EU AI Act, DSGVO Art. 48, NIS2, DORA — "
        "Compliance-Welle als Designprinzip.",
        font=FONT_DISPLAY,
        size=11,
        color=BODY,
        line_spacing=1.4,
        space_before=4,
    )
    add_paragraph(
        body_tb.text_frame,
        "▸  Akt III — Der Spielzug. Nur ein Drittel Ihrer KI muss tatsächlich "
        "souverän sein. Welches Drittel — das ist die einzige Entscheidung.",
        font=FONT_DISPLAY,
        size=11,
        color=BODY,
        line_spacing=1.4,
        space_before=4,
    )

    # Rechte Spalte: 3 Stats vertikal
    stats_x = Inches(9.0)
    stats_w = Inches(3.6)
    stat_y = Inches(2.85)

    for value, label in [
        ("1/3", "Ihrer KI-Workloads muss souverän sein. Der Rest darf bei den globalen Hyperscalern bleiben."),
        ("100  Mio. EUR", "Bußgeld gegen Yango im Mai 2026 wegen unerlaubter Datentransfers."),
        ("30.000", "Unternehmen in Deutschland fallen seit Dezember 2025 unter die NIS2-Pflichten."),
    ]:
        # Value
        add_text(
            slide,
            value,
            stats_x,
            stat_y,
            stats_w,
            Inches(0.7),
            font=FONT_DISPLAY,
            size=34,
            bold=True,
            color=PETROL,
            line_spacing=1,
            spacing=-0.3,
        )
        # Label
        add_text(
            slide,
            label,
            stats_x,
            stat_y + Inches(0.6),
            stats_w,
            Inches(0.7),
            font=FONT_DISPLAY,
            size=10,
            color=BODY_MUTED,
            line_spacing=1.4,
        )
        stat_y += Inches(1.25)

    return slide


def build_act_divider(prs, act_num: str, act_title: str, kicker: str):
    """Section-Divider zwischen den drei Akten."""
    slide = add_blank_slide(prs)

    # Riesige Akt-Nummer (Cousine)
    add_text(
        slide,
        f"AKT {act_num}",
        Inches(0.7),
        Inches(2.5),
        Inches(11.9),
        Inches(0.8),
        font=FONT_MONO,
        size=22,
        color=PETROL,
        spacing=0.50,
        align=PP_ALIGN.CENTER,
    )

    # Akt-Titel
    add_text(
        slide,
        act_title,
        Inches(0.7),
        Inches(3.3),
        Inches(11.9),
        Inches(1.2),
        font=FONT_DISPLAY,
        size=64,
        bold=True,
        color=HEADLINE,
        spacing=-0.6,
        align=PP_ALIGN.CENTER,
        line_spacing=1.05,
    )

    # Petrol-Trennlinie
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.16),
        Inches(4.7),
        Inches(1),
        Inches(0.04),
    )
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = PETROL

    # Kicker
    add_text(
        slide,
        kicker,
        Inches(1.5),
        Inches(5),
        Inches(10.3),
        Inches(1.2),
        font=FONT_DISPLAY,
        size=18,
        color=BODY,
        line_spacing=1.5,
        align=PP_ALIGN.CENTER,
        italic=True,
    )

    return slide


def build_editorial(
    prs,
    eyebrow: str,
    headline: str,
    body_blocks: list[tuple[str, dict]],
    *,
    sidebar: tuple[str, list[tuple[str, str]]] | None = None,
):
    """
    Editorial-Body-Slide.

    body_blocks: Liste (text, kwargs) für add_paragraph
    sidebar: optional (sidebar_title, [(label, body), ...]) für Faktenkarten rechts
    """
    slide = add_blank_slide(prs)

    # Eyebrow
    add_text(
        slide,
        eyebrow,
        Inches(0.7),
        Inches(0.9),
        Inches(11.9),
        Inches(0.4),
        font=FONT_MONO,
        size=10,
        color=PETROL,
        spacing=0.30,
    )

    # H1
    add_text(
        slide,
        headline,
        Inches(0.7),
        Inches(1.3),
        Inches(11.9),
        Inches(1),
        font=FONT_DISPLAY,
        size=30,
        bold=True,
        color=HEADLINE,
        spacing=-0.3,
        line_spacing=1.15,
    )

    hr_line(slide, Inches(2.4))

    # Body
    body_w = Inches(7.5) if sidebar else Inches(11.9)
    body_tb = slide.shapes.add_textbox(
        Inches(0.7), Inches(2.7), body_w, Inches(4.2)
    )
    first = True
    for text, kw in body_blocks:
        if first:
            set_text(
                body_tb,
                text,
                font=kw.get("font", FONT_DISPLAY),
                size=kw.get("size", 13),
                bold=kw.get("bold", False),
                italic=kw.get("italic", False),
                color=kw.get("color", BODY),
                line_spacing=kw.get("line_spacing", 1.5),
            )
            first = False
        else:
            add_paragraph(
                body_tb.text_frame,
                text,
                font=kw.get("font", FONT_DISPLAY),
                size=kw.get("size", 13),
                bold=kw.get("bold", False),
                italic=kw.get("italic", False),
                color=kw.get("color", BODY),
                line_spacing=kw.get("line_spacing", 1.5),
                space_before=kw.get("space_before", 8),
            )

    # Sidebar (rechte Spalte mit Mini-Karten)
    if sidebar:
        sb_title, items = sidebar
        sb_x = Inches(8.7)
        sb_w = Inches(3.9)

        # Sidebar-Title (Mono-Eyebrow)
        add_text(
            slide,
            sb_title,
            sb_x,
            Inches(2.7),
            sb_w,
            Inches(0.4),
            font=FONT_MONO,
            size=9,
            color=PETROL,
            spacing=0.30,
            bold=True,
        )

        # Petrol-Linie unter Title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, sb_x, Inches(3.05), Inches(0.6), Inches(0.015)
        )
        line.line.fill.background()
        line.fill.solid()
        line.fill.fore_color.rgb = PETROL

        # Items
        item_y = Inches(3.2)
        for label, body in items:
            add_text(
                slide,
                label,
                sb_x,
                item_y,
                sb_w,
                Inches(0.3),
                font=FONT_DISPLAY,
                size=11,
                bold=True,
                color=HEADLINE,
            )
            add_text(
                slide,
                body,
                sb_x,
                item_y + Inches(0.3),
                sb_w,
                Inches(1),
                font=FONT_DISPLAY,
                size=9.5,
                color=BODY_MUTED,
                line_spacing=1.4,
            )
            item_y += Inches(1.2)

    return slide


def build_pull_quote(
    prs, eyebrow: str, quote: str, attribution: str, context_blocks: list[tuple[str, str]]
):
    """
    Pull-Quote-Slide für das Carniaux-Zitat.

    context_blocks: [(title, body), ...] — zwei Spalten unter dem Quote
    """
    slide = add_blank_slide(prs)

    # Eyebrow
    add_text(
        slide,
        eyebrow,
        Inches(0.7),
        Inches(0.9),
        Inches(11.9),
        Inches(0.4),
        font=FONT_MONO,
        size=10,
        color=PETROL,
        spacing=0.30,
    )

    # Quote zentriert
    add_text(
        slide,
        f'„{quote}"',
        Inches(1.5),
        Inches(1.6),
        Inches(10.3),
        Inches(1.4),
        font=FONT_DISPLAY,
        size=34,
        bold=True,
        italic=True,
        color=HEADLINE,
        align=PP_ALIGN.CENTER,
        spacing=-0.3,
        line_spacing=1.15,
    )

    # Attribution
    add_text(
        slide,
        f"— {attribution}",
        Inches(1.5),
        Inches(3.15),
        Inches(10.3),
        Inches(0.4),
        font=FONT_MONO,
        size=10,
        color=PETROL,
        spacing=0.20,
        align=PP_ALIGN.CENTER,
    )

    hr_line(slide, Inches(3.85))

    # Zwei-Spalten-Kontext darunter
    col_w = Inches(5.7)
    col_y = Inches(4.1)
    for i, (title, body) in enumerate(context_blocks[:2]):
        col_x = Inches(0.7) + i * (col_w + Inches(0.5))

        add_text(
            slide,
            title,
            col_x,
            col_y,
            col_w,
            Inches(0.35),
            font=FONT_MONO,
            size=10,
            bold=True,
            color=PETROL,
            spacing=0.25,
        )
        add_text(
            slide,
            body,
            col_x,
            col_y + Inches(0.4),
            col_w,
            Inches(2.4),
            font=FONT_DISPLAY,
            size=11,
            color=BODY,
            line_spacing=1.5,
        )

    return slide


def build_four_vignettes(
    prs, eyebrow: str, headline: str, intro: str, vignettes: list[dict]
):
    """4er-Vignetten-Grid (2x2)."""
    slide = add_blank_slide(prs)

    # Eyebrow
    add_text(
        slide,
        eyebrow,
        Inches(0.7),
        Inches(0.9),
        Inches(11.9),
        Inches(0.4),
        font=FONT_MONO,
        size=10,
        color=PETROL,
        spacing=0.30,
    )

    # H1
    add_text(
        slide,
        headline,
        Inches(0.7),
        Inches(1.3),
        Inches(11.9),
        Inches(0.7),
        font=FONT_DISPLAY,
        size=28,
        bold=True,
        color=HEADLINE,
        spacing=-0.3,
    )

    # Intro
    add_text(
        slide,
        intro,
        Inches(0.7),
        Inches(2.05),
        Inches(11.9),
        Inches(0.5),
        font=FONT_DISPLAY,
        size=12,
        italic=True,
        color=BODY_MUTED,
        line_spacing=1.4,
    )

    # 2x2 Grid
    card_w = Inches(5.7)
    card_h = Inches(2.0)
    gap_x = Inches(0.5)
    gap_y = Inches(0.25)
    start_x = Inches(0.7)
    start_y = Inches(2.85)

    for i, vig in enumerate(vignettes[:4]):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        # Petrol-Top-Line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, Inches(0.5), Inches(0.025)
        )
        line.line.fill.background()
        line.fill.solid()
        line.fill.fore_color.rgb = PETROL

        # Eyebrow (Branche)
        add_text(
            slide,
            vig["eyebrow"],
            x,
            y + Inches(0.1),
            card_w,
            Inches(0.3),
            font=FONT_MONO,
            size=9,
            color=PETROL,
            spacing=0.25,
            bold=True,
        )

        # Mini-H2
        add_text(
            slide,
            vig["title"],
            x,
            y + Inches(0.45),
            card_w,
            Inches(0.5),
            font=FONT_DISPLAY,
            size=15,
            bold=True,
            color=HEADLINE,
            line_spacing=1.2,
        )

        # Body
        add_text(
            slide,
            vig["body"],
            x,
            y + Inches(0.95),
            card_w,
            card_h - Inches(0.95),
            font=FONT_DISPLAY,
            size=9.5,
            color=BODY,
            line_spacing=1.45,
        )

    return slide


def build_triade(prs, eyebrow: str, headline: str, intro: str, items: list[dict]):
    """Drei nummerierte Kacheln horizontal — z. B. Agency-Triade."""
    slide = add_blank_slide(prs)

    add_text(
        slide,
        eyebrow,
        Inches(0.7),
        Inches(0.9),
        Inches(11.9),
        Inches(0.4),
        font=FONT_MONO,
        size=10,
        color=PETROL,
        spacing=0.30,
    )

    add_text(
        slide,
        headline,
        Inches(0.7),
        Inches(1.3),
        Inches(11.9),
        Inches(0.8),
        font=FONT_DISPLAY,
        size=30,
        bold=True,
        color=HEADLINE,
        spacing=-0.3,
        line_spacing=1.15,
    )

    add_text(
        slide,
        intro,
        Inches(0.7),
        Inches(2.2),
        Inches(11.9),
        Inches(1.4),
        font=FONT_DISPLAY,
        size=12.5,
        color=BODY,
        line_spacing=1.5,
        italic=True,
    )

    # Drei Kacheln
    card_w = Inches(3.8)
    gap = Inches(0.3)
    total_w = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) / 2
    card_y = Inches(4.2)
    card_h = Inches(2.5)

    for i, item in enumerate(items[:3]):
        x = start_x + i * (card_w + gap)

        # Top Petrol-Linie
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, card_y, card_w, Inches(0.025)
        )
        line.line.fill.background()
        line.fill.solid()
        line.fill.fore_color.rgb = PETROL

        # Num
        add_text(
            slide,
            item["num"],
            x,
            card_y + Inches(0.1),
            card_w,
            Inches(0.35),
            font=FONT_MONO,
            size=11,
            color=PETROL,
            spacing=0.20,
        )

        # Title
        add_text(
            slide,
            item["title"],
            x,
            card_y + Inches(0.45),
            card_w,
            Inches(0.5),
            font=FONT_DISPLAY,
            size=17,
            bold=True,
            color=HEADLINE,
            line_spacing=1.15,
        )

        # Body
        add_text(
            slide,
            item["body"],
            x,
            card_y + Inches(0.95),
            card_w,
            card_h - Inches(0.95),
            font=FONT_DISPLAY,
            size=10.5,
            color=BODY,
            line_spacing=1.5,
        )

    return slide


def build_two_col(prs, eyebrow: str, headline: str, intro: str, left: dict, right: dict):
    """Zwei gleichwertige Spalten — Intelligence Taker vs Maker."""
    slide = add_blank_slide(prs)

    add_text(
        slide,
        eyebrow,
        Inches(0.7),
        Inches(0.9),
        Inches(11.9),
        Inches(0.4),
        font=FONT_MONO,
        size=10,
        color=PETROL,
        spacing=0.30,
    )

    add_text(
        slide,
        headline,
        Inches(0.7),
        Inches(1.3),
        Inches(11.9),
        Inches(0.8),
        font=FONT_DISPLAY,
        size=30,
        bold=True,
        color=HEADLINE,
        spacing=-0.3,
        line_spacing=1.15,
    )

    if intro:
        add_text(
            slide,
            intro,
            Inches(0.7),
            Inches(2.15),
            Inches(11.9),
            Inches(0.7),
            font=FONT_DISPLAY,
            size=12,
            color=BODY,
            line_spacing=1.5,
            italic=True,
        )
        col_y = Inches(3.05)
    else:
        col_y = Inches(2.3)

    hr_line(slide, col_y - Inches(0.15))

    col_w = Inches(5.7)
    col_h = Inches(3.8)

    for i, col in enumerate([left, right]):
        x = Inches(0.7) + i * (col_w + Inches(0.5))
        accent = col.get("accent", PETROL)

        # Petrol-Top-Line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, col_y, Inches(0.5), Inches(0.03)
        )
        line.line.fill.background()
        line.fill.solid()
        line.fill.fore_color.rgb = accent

        # Eyebrow
        add_text(
            slide,
            col["eyebrow"],
            x,
            col_y + Inches(0.1),
            col_w,
            Inches(0.3),
            font=FONT_MONO,
            size=10,
            color=accent,
            spacing=0.25,
            bold=True,
        )

        # Title
        add_text(
            slide,
            col["title"],
            x,
            col_y + Inches(0.45),
            col_w,
            Inches(0.5),
            font=FONT_DISPLAY,
            size=20,
            bold=True,
            color=HEADLINE,
        )

        # Bullets
        bullet_tb = slide.shapes.add_textbox(
            x, col_y + Inches(1.05), col_w, col_h - Inches(1.05)
        )
        first = True
        for bp in col["bullets"]:
            if first:
                set_text(
                    bullet_tb,
                    f"▸  {bp}",
                    font=FONT_DISPLAY,
                    size=10.5,
                    color=BODY,
                    line_spacing=1.5,
                )
                first = False
            else:
                add_paragraph(
                    bullet_tb.text_frame,
                    f"▸  {bp}",
                    font=FONT_DISPLAY,
                    size=10.5,
                    color=BODY,
                    line_spacing=1.5,
                    space_before=6,
                )

    return slide


def build_four_models(prs, eyebrow: str, headline: str, intro: str, models: list[dict]):
    """Vier Modell-Karten horizontal (Sovereign Operating Matrix)."""
    slide = add_blank_slide(prs)

    add_text(
        slide,
        eyebrow,
        Inches(0.7),
        Inches(0.9),
        Inches(11.9),
        Inches(0.4),
        font=FONT_MONO,
        size=10,
        color=PETROL,
        spacing=0.30,
    )

    add_text(
        slide,
        headline,
        Inches(0.7),
        Inches(1.3),
        Inches(11.9),
        Inches(0.7),
        font=FONT_DISPLAY,
        size=28,
        bold=True,
        color=HEADLINE,
        spacing=-0.3,
    )

    if intro:
        add_text(
            slide,
            intro,
            Inches(0.7),
            Inches(2.05),
            Inches(11.9),
            Inches(0.5),
            font=FONT_DISPLAY,
            size=11.5,
            italic=True,
            color=BODY_MUTED,
            line_spacing=1.4,
        )

    card_w = Inches(2.85)
    gap = Inches(0.15)
    start_x = Inches(0.7)
    card_y = Inches(2.85)
    card_h = Inches(4.0)

    for i, model in enumerate(models[:4]):
        x = start_x + i * (card_w + gap)

        # Top-Linie
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, card_y, card_w, Inches(0.025)
        )
        line.line.fill.background()
        line.fill.solid()
        line.fill.fore_color.rgb = PETROL

        # Num
        add_text(
            slide,
            f"MODELL {i + 1}",
            x,
            card_y + Inches(0.1),
            card_w,
            Inches(0.3),
            font=FONT_MONO,
            size=9,
            color=PETROL,
            spacing=0.25,
            bold=True,
        )

        # Name
        add_text(
            slide,
            model["name"],
            x,
            card_y + Inches(0.4),
            card_w,
            Inches(0.7),
            font=FONT_DISPLAY,
            size=14,
            bold=True,
            color=HEADLINE,
            line_spacing=1.15,
        )

        # Tagline
        add_text(
            slide,
            model["tagline"],
            x,
            card_y + Inches(1.05),
            card_w,
            Inches(0.5),
            font=FONT_DISPLAY,
            size=9.5,
            color=BODY,
            italic=True,
            line_spacing=1.35,
        )

        # Bullets
        bullet_tb = slide.shapes.add_textbox(
            x, card_y + Inches(1.55), card_w, Inches(1.8)
        )
        first = True
        for bp in model["bullets"]:
            if first:
                set_text(
                    bullet_tb,
                    f"·  {bp}",
                    font=FONT_DISPLAY,
                    size=9,
                    color=BODY,
                    line_spacing=1.4,
                )
                first = False
            else:
                add_paragraph(
                    bullet_tb.text_frame,
                    f"·  {bp}",
                    font=FONT_DISPLAY,
                    size=9,
                    color=BODY,
                    line_spacing=1.4,
                )

        # Für
        add_text(
            slide,
            f"FÜR  ·  {model['fuer']}",
            x,
            card_y + Inches(3.45),
            card_w,
            Inches(0.5),
            font=FONT_MONO,
            size=8,
            color=PETROL_DEEP,
            spacing=0.2,
            line_spacing=1.35,
        )

    return slide


def build_csd_framework(prs):
    """CSD-Framework: Control / Steer / Depend mit gemeinsamem Footer."""
    slide = add_blank_slide(prs)

    add_text(
        slide,
        "AKT III — DER SPIELZUG  ·  SZENE 1",
        Inches(0.7),
        Inches(0.9),
        Inches(11.9),
        Inches(0.4),
        font=FONT_MONO,
        size=10,
        color=PETROL,
        spacing=0.30,
    )

    add_text(
        slide,
        "Das CSD-Framework: Control · Steer · Depend.",
        Inches(0.7),
        Inches(1.3),
        Inches(11.9),
        Inches(0.7),
        font=FONT_DISPLAY,
        size=28,
        bold=True,
        color=HEADLINE,
        spacing=-0.3,
    )

    add_text(
        slide,
        "Sie müssen nicht alles selbst bauen. Sie müssen nur die richtigen Dinge "
        "selbst kontrollieren. Das CSD-Framework sortiert KI-Workloads in drei "
        "Kategorien — und macht damit die einzige strategische Entscheidung "
        "explizit, die heute zählt.",
        Inches(0.7),
        Inches(2.05),
        Inches(11.9),
        Inches(0.7),
        font=FONT_DISPLAY,
        size=11.5,
        italic=True,
        color=BODY_MUTED,
        line_spacing=1.4,
    )

    csd = [
        {
            "letter": "C",
            "name": "Control",
            "share": "≈ 10–30 %",
            "tagline": "Missionskritische Intelligenz",
            "examples": "Wirkstoffforschung, autonome Produktionssteuerung, Kern-Bankensysteme, Patent-Analysen.",
            "consequence": "Souveräne Infrastruktur (Managed Sovereignty / Bare Metal in EU). Open-Weight-Modelle. BYOK. Volle Auditierbarkeit.",
            "accent": ROST,
        },
        {
            "letter": "S",
            "name": "Steer",
            "share": "≈ 30–60 %",
            "tagline": "Wichtige Prozesse mit Gestaltungsanspruch",
            "examples": "HR-Systeme, Support-Bots, CRM-Analysen, Logistikoptimierung, Marketing-Automation.",
            "consequence": "Souveränes MaaS bei zertifizierten Providern. Compliance-Wrappers. Klare Datenverträge.",
            "accent": BERNSTEIN,
        },
        {
            "letter": "D",
            "name": "Depend",
            "share": "≈ 10–60 %",
            "tagline": "Bewusste Abhängigkeit für unkritische Dienste",
            "examples": "Coding-Assistenten, Marketing-Texte, Brainstorming, Übersetzungen, generische Bildgenerierung.",
            "consequence": "Standard-APIs (OpenAI, Anthropic, Google). Datenklassifizierung verhindert sensiblen Daten-Abfluss.",
            "accent": SALBEI,
        },
    ]

    card_w = Inches(3.85)
    gap = Inches(0.2)
    start_x = Inches(0.7)
    card_y = Inches(2.95)
    card_h = Inches(3.4)

    for i, c in enumerate(csd):
        x = start_x + i * (card_w + gap)

        # Top-Linie in Akzent-Farbe
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, card_y, card_w, Inches(0.035)
        )
        line.line.fill.background()
        line.fill.solid()
        line.fill.fore_color.rgb = c["accent"]

        # Big Letter
        add_text(
            slide,
            c["letter"],
            x,
            card_y + Inches(0.15),
            Inches(0.7),
            Inches(0.9),
            font=FONT_DISPLAY,
            size=56,
            bold=True,
            color=c["accent"],
            line_spacing=1,
            spacing=-0.6,
        )

        # Name + Share (rechts vom Letter)
        add_text(
            slide,
            c["name"],
            x + Inches(0.85),
            card_y + Inches(0.3),
            card_w - Inches(0.85),
            Inches(0.4),
            font=FONT_DISPLAY,
            size=18,
            bold=True,
            color=HEADLINE,
        )
        add_text(
            slide,
            f"{c['share']} der Workloads",
            x + Inches(0.85),
            card_y + Inches(0.72),
            card_w - Inches(0.85),
            Inches(0.3),
            font=FONT_MONO,
            size=8.5,
            color=BODY_MUTED,
            spacing=0.18,
        )

        # Tagline
        add_text(
            slide,
            c["tagline"],
            x,
            card_y + Inches(1.2),
            card_w,
            Inches(0.4),
            font=FONT_DISPLAY,
            size=11,
            italic=True,
            color=PETROL_DEEP,
            line_spacing=1.3,
        )

        # Beispiele
        add_text(
            slide,
            "BEISPIELE",
            x,
            card_y + Inches(1.7),
            card_w,
            Inches(0.25),
            font=FONT_MONO,
            size=7.5,
            color=PETROL,
            spacing=0.25,
            bold=True,
        )
        add_text(
            slide,
            c["examples"],
            x,
            card_y + Inches(1.95),
            card_w,
            Inches(0.7),
            font=FONT_DISPLAY,
            size=9.5,
            color=BODY,
            line_spacing=1.4,
        )

        # Konsequenz
        add_text(
            slide,
            "KONSEQUENZ",
            x,
            card_y + Inches(2.5),
            card_w,
            Inches(0.25),
            font=FONT_MONO,
            size=7.5,
            color=PETROL,
            spacing=0.25,
            bold=True,
        )
        add_text(
            slide,
            c["consequence"],
            x,
            card_y + Inches(2.75),
            card_w,
            Inches(0.7),
            font=FONT_DISPLAY,
            size=9.5,
            color=BODY,
            line_spacing=1.4,
        )

    # Footer-Statement
    add_text(
        slide,
        "Die einzige strategische Entscheidung, die heute zählt:",
        Inches(0.7),
        Inches(6.6),
        Inches(11.9),
        Inches(0.4),
        font=FONT_MONO,
        size=10,
        color=PETROL,
        spacing=0.20,
        align=PP_ALIGN.CENTER,
        bold=True,
    )
    add_text(
        slide,
        "Welches Drittel Ihrer KI gehört in welche Kategorie?",
        Inches(0.7),
        Inches(6.95),
        Inches(11.9),
        Inches(0.4),
        font=FONT_DISPLAY,
        size=17,
        bold=True,
        color=HEADLINE,
        align=PP_ALIGN.CENTER,
    )

    return slide


def build_roadmap(prs):
    """90-Tage-Roadmap: drei Phasen-Spalten mit Schritten."""
    slide = add_blank_slide(prs)

    add_text(
        slide,
        "AKT III — DER SPIELZUG  ·  SZENE 2",
        Inches(0.7),
        Inches(0.9),
        Inches(11.9),
        Inches(0.4),
        font=FONT_MONO,
        size=10,
        color=PETROL,
        spacing=0.30,
    )

    add_text(
        slide,
        "Die 90-Tage-Roadmap.",
        Inches(0.7),
        Inches(1.3),
        Inches(11.9),
        Inches(0.8),
        font=FONT_DISPLAY,
        size=30,
        bold=True,
        color=HEADLINE,
        spacing=-0.3,
    )

    add_text(
        slide,
        "Souveränität ist kein Big-Bang-Projekt. Sieben Schritte in drei Phasen — "
        "pragmatisch, sequenziell, ohne Risiko-Sprünge.",
        Inches(0.7),
        Inches(2.1),
        Inches(11.9),
        Inches(0.5),
        font=FONT_DISPLAY,
        size=12,
        italic=True,
        color=BODY_MUTED,
        line_spacing=1.4,
    )

    phases = [
        {
            "label": "PHASE 1 · TAG 1–30",
            "title": "Inventur & Klassifizierung",
            "steps": [
                ("1.  KI-Inventur", "Vollständige Erfassung aller KI-Anwendungen — inklusive Schatten-KI ohne IT-Abstimmung."),
                ("2.  Risiko-Klassifizierung", "Nach EU AI Act einordnen: verboten / hoch / begrenzt / minimal. Auch Art. 6 Abs. 3."),
                ("3.  CSD-Einordnung", "Jedes System in Control / Steer / Depend einordnen — wird zum Steuerungs-Cockpit."),
            ],
            "accent": ROST,
        },
        {
            "label": "PHASE 2 · TAG 30–60",
            "title": "Risikoanalyse & Strategie",
            "steps": [
                ("4.  DORA-Konzentrationsrisiko", "Welche Anbieter halten kritische Funktionen? Wo wäre Wechsel in 6 Monaten unmöglich?"),
                ("5.  Datenklassifizierung", "Drei Stufen: öffentlich / geschäftlich vertraulich / regulatorisch kritisch. Nur Stufe 3 zwingend souverän."),
                ("6.  Exit-Strategie", "Pro Control-Anbieter ein konkreter, dokumentierter Wechsel-Plan."),
            ],
            "accent": BERNSTEIN,
        },
        {
            "label": "PHASE 3 · TAG 60–90",
            "title": "Erste Migration & Governance",
            "steps": [
                ("7.  Pilot-Migration", "Eine Control-Workload auf Open-Weight-Modell in souveräner Umgebung (Mistral auf STACKIT / IONOS). Governance + Schulung parallel."),
            ],
            "accent": SALBEI,
        },
    ]

    col_w = Inches(3.95)
    gap = Inches(0.15)
    start_x = Inches(0.7)
    col_y = Inches(2.85)

    for i, ph in enumerate(phases):
        x = start_x + i * (col_w + gap)

        # Top-Linie
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, col_y, col_w, Inches(0.03)
        )
        line.line.fill.background()
        line.fill.solid()
        line.fill.fore_color.rgb = ph["accent"]

        # Phase-Label
        add_text(
            slide,
            ph["label"],
            x,
            col_y + Inches(0.1),
            col_w,
            Inches(0.3),
            font=FONT_MONO,
            size=9,
            color=ph["accent"],
            spacing=0.25,
            bold=True,
        )

        # Title
        add_text(
            slide,
            ph["title"],
            x,
            col_y + Inches(0.4),
            col_w,
            Inches(0.5),
            font=FONT_DISPLAY,
            size=16,
            bold=True,
            color=HEADLINE,
        )

        # Steps
        step_y = col_y + Inches(1.05)
        for step_label, step_body in ph["steps"]:
            add_text(
                slide,
                step_label,
                x,
                step_y,
                col_w,
                Inches(0.3),
                font=FONT_DISPLAY,
                size=11,
                bold=True,
                color=HEADLINE,
            )
            add_text(
                slide,
                step_body,
                x,
                step_y + Inches(0.3),
                col_w,
                Inches(0.85),
                font=FONT_DISPLAY,
                size=9.5,
                color=BODY,
                line_spacing=1.45,
            )
            step_y += Inches(1.2)

    return slide


def build_check_questions(prs):
    """12-Fragen-Souveränitäts-Check als 2×2-Grid."""
    slide = add_blank_slide(prs)

    add_text(
        slide,
        "AKT III — DER SPIELZUG  ·  SZENE 3",
        Inches(0.7),
        Inches(0.9),
        Inches(11.9),
        Inches(0.4),
        font=FONT_MONO,
        size=10,
        color=PETROL,
        spacing=0.30,
    )

    add_text(
        slide,
        "Der 12-Fragen-Souveränitäts-Check.",
        Inches(0.7),
        Inches(1.3),
        Inches(11.9),
        Inches(0.7),
        font=FONT_DISPLAY,
        size=28,
        bold=True,
        color=HEADLINE,
        spacing=-0.3,
    )

    add_text(
        slide,
        'Ein ehrlicher Spiegel für die Geschäftsführung. Beantworten Sie jede '
        'Frage mit Ja oder Nein. „Ich weiß es nicht" zählt als Nein.',
        Inches(0.7),
        Inches(2.05),
        Inches(11.9),
        Inches(0.5),
        font=FONT_DISPLAY,
        size=11.5,
        italic=True,
        color=BODY_MUTED,
        line_spacing=1.4,
    )

    cats = [
        {
            "name": "KATEGORIE I  ·  Daten-Kontrolle",
            "qs": [
                "Wissen Sie exakt, in welchen Jurisdiktionen Ihre KI-Trainingsdaten gespeichert und verarbeitet werden?",
                "Können Sie garantieren, dass kritische Daten niemals Grenzen verlassen, die Sie nicht freigegeben haben?",
                "Halten Sie die Verschlüsselungs-Keys (BYOK) selbst — nicht der Provider?",
            ],
        },
        {
            "name": "KATEGORIE II  ·  Modell-Kontrolle",
            "qs": [
                "Besitzen Sie die Eigentumsrechte an den KI-Modellen — oder sind Sie nur Lizenznehmer einer proprietären API?",
                "Können Sie Ihre Modelle unabhängig modifizieren, verbessern oder auditieren (Open Weights)?",
                "Funktioniert Ihre KI weiter, wenn die Geschäftsbeziehung zum Vendor morgen endet?",
            ],
        },
        {
            "name": "KATEGORIE III  ·  Infrastruktur-Kontrolle",
            "qs": [
                "Könnten Ihre KI-Systeme weiterlaufen, wenn die Verbindung zu einem ausländischen Provider unterbrochen wäre?",
                "Wissen Sie genau, wer administrativen Zugriff auf die Infrastruktur hat?",
                "Sind Ihre KI-Systeme frei von Abhängigkeiten, die unter Exportkontrolle oder fremder Jurisdiktion stehen?",
            ],
        },
        {
            "name": "KATEGORIE IV  ·  Entscheidungs-Kontrolle",
            "qs": [
                "Können Sie (oder ein unabhängiger Dritter) erklären, wie Ihre KI zu einer bestimmten Entscheidung kam?",
                "Haben Sie die technische Möglichkeit, KI-Verhalten an neue Regulatorik anzupassen (z. B. EU AI Act)?",
                "Sind Ihre KI-Fähigkeiten immun gegen einseitige Richtlinien-Änderungen eines Drittanbieters?",
            ],
        },
    ]

    card_w = Inches(5.85)
    card_h = Inches(2.0)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)
    start_x = Inches(0.7)
    start_y = Inches(2.8)

    for i, cat in enumerate(cats):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        # Cat-Label
        add_text(
            slide,
            cat["name"],
            x,
            y,
            card_w,
            Inches(0.3),
            font=FONT_MONO,
            size=9.5,
            color=PETROL,
            spacing=0.25,
            bold=True,
        )

        # Petrol-Underline
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y + Inches(0.32), Inches(0.5), Inches(0.018)
        )
        line.line.fill.background()
        line.fill.solid()
        line.fill.fore_color.rgb = PETROL

        # Questions
        q_tb = slide.shapes.add_textbox(
            x, y + Inches(0.45), card_w, card_h - Inches(0.45)
        )
        first = True
        for q in cat["qs"]:
            if first:
                set_text(
                    q_tb,
                    f"☐  {q}",
                    font=FONT_DISPLAY,
                    size=9.5,
                    color=BODY,
                    line_spacing=1.4,
                )
                first = False
            else:
                add_paragraph(
                    q_tb.text_frame,
                    f"☐  {q}",
                    font=FONT_DISPLAY,
                    size=9.5,
                    color=BODY,
                    line_spacing=1.4,
                    space_before=4,
                )

    return slide


def build_levels(prs):
    """Score-Levels nach 12-Fragen-Check: 4 Levels horizontal."""
    slide = add_blank_slide(prs)

    add_text(
        slide,
        "AUSWERTUNG  ·  12-FRAGEN-CHECK",
        Inches(0.7),
        Inches(0.9),
        Inches(11.9),
        Inches(0.4),
        font=FONT_MONO,
        size=10,
        color=PETROL,
        spacing=0.30,
    )

    add_text(
        slide,
        "Vier Souveränitäts-Level.",
        Inches(0.7),
        Inches(1.3),
        Inches(11.9),
        Inches(0.8),
        font=FONT_DISPLAY,
        size=30,
        bold=True,
        color=HEADLINE,
        spacing=-0.3,
    )

    add_text(
        slide,
        "Zählen Sie die Ja-Antworten aus dem 12-Fragen-Check. Das Ergebnis "
        "verortet Ihr Unternehmen auf einem von vier Souveränitäts-Levels.",
        Inches(0.7),
        Inches(2.1),
        Inches(11.9),
        Inches(0.5),
        font=FONT_DISPLAY,
        size=12,
        italic=True,
        color=BODY_MUTED,
        line_spacing=1.4,
    )

    levels = [
        {"score": "0–3", "name": "Level 1", "title": "Volle Abhängigkeit", "body": "Sie sind reiner Intelligence Taker. Maximales Lock-in- und Compliance-Risiko.", "accent": ROST},
        {"score": "4–6", "name": "Level 2", "title": "Hybride Kontrolle", "body": "Erste Schritte gemacht — Kern-Intelligenz bleibt aber in fremden Händen.", "accent": BERNSTEIN},
        {"score": "7–9", "name": "Level 3", "title": "Gesteuerte Souveränität", "body": "Sie nutzen dedizierte Infrastruktur, erfüllen Compliance proaktiv.", "accent": PETROL},
        {"score": "10–12", "name": "Level 4", "title": "Komplette Souveränität", "body": "Sie sind Intelligence Maker. KI ist strategisches Asset.", "accent": SALBEI},
    ]

    card_w = Inches(2.85)
    gap = Inches(0.15)
    start_x = Inches(0.7)
    card_y = Inches(2.85)
    card_h = Inches(3.2)

    for i, lvl in enumerate(levels):
        x = start_x + i * (card_w + gap)

        # Top-Linie
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, card_y, card_w, Inches(0.035)
        )
        line.line.fill.background()
        line.fill.solid()
        line.fill.fore_color.rgb = lvl["accent"]

        add_text(
            slide,
            lvl["score"] + "  JA",
            x,
            card_y + Inches(0.15),
            card_w,
            Inches(0.4),
            font=FONT_MONO,
            size=11,
            color=lvl["accent"],
            spacing=0.25,
            bold=True,
        )

        add_text(
            slide,
            lvl["name"],
            x,
            card_y + Inches(0.55),
            card_w,
            Inches(0.5),
            font=FONT_DISPLAY,
            size=22,
            bold=True,
            color=HEADLINE,
            spacing=-0.2,
        )

        add_text(
            slide,
            lvl["title"],
            x,
            card_y + Inches(1.1),
            card_w,
            Inches(0.5),
            font=FONT_DISPLAY,
            size=14,
            italic=True,
            color=lvl["accent"],
            line_spacing=1.2,
        )

        add_text(
            slide,
            lvl["body"],
            x,
            card_y + Inches(1.7),
            card_w,
            card_h - Inches(1.7),
            font=FONT_DISPLAY,
            size=10,
            color=BODY,
            line_spacing=1.5,
        )

    return slide


def build_closing(prs):
    """Schlussbild: drei Sätze für die nächste Vorstandssitzung."""
    slide = add_blank_slide(prs)

    add_text(
        slide,
        "AKT III · SCHLUSSBILD",
        Inches(0.7),
        Inches(0.9),
        Inches(11.9),
        Inches(0.4),
        font=FONT_MONO,
        size=10,
        color=PETROL,
        spacing=0.30,
    )

    add_text(
        slide,
        "Auf eigenem Grund bauen.",
        Inches(0.7),
        Inches(1.4),
        Inches(11.9),
        Inches(1),
        font=FONT_DISPLAY,
        size=42,
        bold=True,
        color=HEADLINE,
        spacing=-0.5,
    )

    add_text(
        slide,
        "DREI SÄTZE FÜR DIE NÄCHSTE VORSTANDSSITZUNG",
        Inches(0.7),
        Inches(2.7),
        Inches(11.9),
        Inches(0.4),
        font=FONT_MONO,
        size=11,
        color=PETROL,
        spacing=0.30,
        bold=True,
    )

    hr_line(slide, Inches(3.15))

    statements = [
        '01  „Wir wissen, welches Drittel unserer KI tatsächlich souverän sein muss — und welches darf bei den Hyperscalern bleiben."',
        '02  „Wir haben für jeden Control-Anbieter eine dokumentierte Exit-Strategie nach DORA-Art."',
        '03  „Wir nutzen die Compliance-Welle 2025–2028 als Designprinzip — nicht als Last."',
    ]

    y = Inches(3.6)
    for s in statements:
        add_text(
            slide,
            s,
            Inches(1),
            y,
            Inches(11.3),
            Inches(0.8),
            font=FONT_DISPLAY,
            size=15,
            color=HEADLINE,
            line_spacing=1.5,
            italic=True,
        )
        y += Inches(1)

    return slide


def build_sources(prs):
    """Quellen-Slide (28 Einträge auf 2 Spalten)."""
    slide = add_blank_slide(prs)

    add_text(
        slide,
        "ANHANG  ·  ALLE QUELLEN MIT URLS  ·  STAND MAI 2026",
        Inches(0.7),
        Inches(0.9),
        Inches(11.9),
        Inches(0.4),
        font=FONT_MONO,
        size=10,
        color=PETROL,
        spacing=0.30,
    )

    add_text(
        slide,
        "Quellen.",
        Inches(0.7),
        Inches(1.3),
        Inches(11.9),
        Inches(0.7),
        font=FONT_DISPLAY,
        size=28,
        bold=True,
        color=HEADLINE,
        spacing=-0.3,
    )

    sources = [
        ("EU AI Act — Hochrisiko-Pflichten und Stichtage", "advisori.de · Mai 2026"),
        ("EU AI Act — Übersicht Europäische Kommission", "digital-strategy.ec.europa.eu"),
        ("EU AI Act — Digital Omnibus, neue Fristen", "ad-hoc-news.de · Mai 2026"),
        ("EU AI Act 2026 — Zwischenstand", "TÜV Consulting · April 2026"),
        ("Microsoft France vor dem Senat (Carniaux, unter Eid)", "heise.de · Juni 2025"),
        ("Microsoft kann US-Zugriff auf EU-Cloud nicht verhindern", "Dr. Datenschutz"),
        ("Souveräne EU-Cloud-Debakel — Senatsanhörung", "Born's IT-Blog · Juli 2025"),
        ("EU Data Act vs. CLOUD Act — Souveränitätskonflikt", "kiteworks.com · März 2026"),
        ("CLOUD Act und europäischer Datenschutz", "kiteworks.com · Februar 2026"),
        ("CLOUD Act-Risiko — Rechtsgutachten Uni Köln (BMI)", "novalnet.de"),
        ("Sicherer Datenaustausch trotz CLOUD Act", "netfiles.com"),
        ("DSGVO-Boundary nichts wert? Microsoft-Eingeständnis", "Datenschutzbeauftragter Hamburg · Aug 2025"),
        ("NIS2 vs. DORA — Unterschiede im Praxis-Guide", "activeMind AG"),
        ("NIS2-Umsetzungsgesetz in Deutschland", "OpenKRITIS"),
        ("NIS2 — Pflichten, Fristen, Checkliste", "SoSafe · 2026"),
        ("NIS-2-Richtlinie — Auswirkungen", "vasgard.com"),
        ("NIS2 trifft DORA — Änderungen Finanzunternehmen", "PayTechLaw"),
        ("DORA, NIS-2 und KRITIS — Anforderungen", "WG-DATA · April 2026"),
        ("Europäische KI-Anbieter und Modelle 2026", "Data Unplugged"),
        ("Mistral AI — Europas Antwort, SAP-Partnerschaft", "Stefan Pfeiffer Blog · März 2026"),
        ("Aleph Alpha, Mistral, LightOn — Europas neue KI-Ära", "iseremo.com"),
        ("Cloud-Souveränität 2026 — Gaia-X, Sovereign Cloud", "cloudmagazin · März 2026"),
        ("Digitale Souveränität 2026 — Delos Cloud, Gaia-X", "digital-chiefs.de · März 2026"),
        ("pluscloud open — Gaia-X-Gründungsmitglied, BSI-C5", "plusserver"),
        ("Cloud Market Share Q3 2025", "QuantumRun Consulting"),
        ("Cloud Market Share 2026 — Top Providers", "Holori"),
        ("Harvard-Studie — Open Source als digitales Fundament", "The Decoder · März 2025"),
        ("Open Source wirtschaftlich 8,8 Billionen USD wert", "heise online"),
    ]

    # 2 Spalten
    col_w = Inches(5.9)
    gap_x = Inches(0.1)
    start_x = Inches(0.7)
    start_y = Inches(2.15)

    per_col = (len(sources) + 1) // 2

    for col in range(2):
        x = start_x + col * (col_w + gap_x)
        tb = slide.shapes.add_textbox(x, start_y, col_w, Inches(4.7))
        items = sources[col * per_col : (col + 1) * per_col]

        first = True
        for i, (title, src) in enumerate(items):
            num = col * per_col + i + 1
            text = f"{num:02d}  {title}"

            if first:
                set_text(
                    tb,
                    text,
                    font=FONT_DISPLAY,
                    size=9,
                    color=HEADLINE,
                    bold=True,
                    line_spacing=1.3,
                )
                first = False
            else:
                add_paragraph(
                    tb.text_frame,
                    text,
                    font=FONT_DISPLAY,
                    size=9,
                    color=HEADLINE,
                    bold=True,
                    line_spacing=1.3,
                    space_before=6,
                )
            add_paragraph(
                tb.text_frame,
                src,
                font=FONT_MONO,
                size=7.5,
                color=MONDSTEIN,
                spacing=0.15,
                line_spacing=1.3,
            )

    return slide


# --------------------------------------------------------------------------
# Main
# --------------------------------------------------------------------------


def main():
    repo_root = Path(__file__).resolve().parent.parent
    output = repo_root / "public" / "whitepaper.pptx"
    output.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- Slide-Reihenfolge ------------------------------------------------

    builders = [
        # Cover + Auftakt
        lambda: build_cover(prs),
        lambda: build_thesis(prs),
        lambda: build_executive_summary(prs),

        # Akt I
        lambda: build_act_divider(
            prs,
            "I",
            "Der Weckruf.",
            "Warum europäische Unternehmen ihre KI heute mehr mieten als besitzen — und was das kostet.",
        ),
        lambda: build_editorial(
            prs,
            "AKT I — DER WECKRUF  ·  SZENE 1",
            "6:14 Uhr. Die Maschinen schweigen.",
            [
                ("Stellen Sie sich folgende Szene vor.", {"size": 14, "italic": True, "color": BODY}),
                (
                    "Ein mittelständischer Automobilzulieferer in Schwaben, 2.000 Mitarbeitende, "
                    "Weltmarktführer in seiner Nische. Die KI-gesteuerte Qualitätskontrolle steht still. "
                    "Kein Defekt — der US-Cloud-Anbieter hat über Nacht die Terms of Service geändert. "
                    "Die API antwortet nicht mehr.",
                    {"size": 12},
                ),
                (
                    "Die Produktionslinie steht. Jede Stunde kostet 180.000 Euro. Der Geschäftsführer "
                    "hört die Stille der Maschinen und begreift zum ersten Mal, was Abhängigkeit "
                    "wirklich bedeutet.",
                    {"size": 12},
                ),
            ],
            sidebar=(
                "DREI FAKTEN ZUR SZENE",
                [
                    ("AGB-Änderung", "Cloud-Anbieter haben Nutzungsbedingungen in der Vergangenheit kurzfristig und einseitig geändert — auch in regulierten Branchen."),
                    ("KI in der Produktion", "KI-gestützte Qualitätskontrolle und Predictive Maintenance sind in der Automobil-Zulieferindustrie Standard — keine Zukunftsmusik."),
                    ("180.000 EUR / Stunde", "Konservative Schätzung für eine durchschnittliche OEM-Tier-1-Produktion bei vollständigem Stillstand."),
                ],
            ),
        ),
        lambda: build_pull_quote(
            prs,
            "AKT I — DER WECKRUF  ·  SZENE 2  ·  Die unsichtbare Kette",
            "Non, je ne peux pas le garantir.",
            "Anton Carniaux, Justiziar Microsoft France  ·  unter Eid, Senatsanhörung 10. Juni 2025",
            [
                (
                    "WAS PASSIERT IST",
                    "Vor dem französischen Senat fragte Berichterstatter Dany Wattebled den "
                    "Microsoft-Justiziar Anton Carniaux unter Eid: Könne Microsoft garantieren, dass "
                    "Daten französischer Bürger niemals ohne Zustimmung der französischen Behörden "
                    "an US-Behörden übermittelt würden? Vier Worte als Antwort. Bestätigt unter Eid, "
                    "was Datenschutzanwälte seit Jahren warnten.",
                ),
                (
                    "WAS ES BEDEUTET",
                    "Der US CLOUD Act gilt für alle US-Anbieter unabhängig vom Speicherort der Daten. "
                    'EU-Rechenzentren, „EU Data Boundary"-Initiativen, vertragliche Zusicherungen — '
                    "sie ändern nichts an der rechtlichen Grundlage. Mai 2026: 100 Mio. EUR Bußgeld "
                    "gegen Yango wegen unerlaubter Datentransfers — die Behörden schöpfen den "
                    "Strafrahmen aus.",
                ),
            ],
        ),
        lambda: build_four_vignettes(
            prs,
            "AKT I — DER WECKRUF  ·  SZENE 3",
            "Vier Branchen, vier Schmerzpunkte.",
            "Die Folgen der Cloud-Abhängigkeit treffen jede regulierte Branche anders. Vier Vignetten.",
            [
                {
                    "eyebrow": "INDUSTRIE  ·  OPERATIVE RESILIENZ",
                    "title": "Der Schwabe um 6:14 Uhr",
                    "body": "Wenn Qualitätskontrolle und Predictive Maintenance in der Cloud eines globalen Anbieters laufen, ist die Produktion an dessen Verfügbarkeit gekoppelt. Drei Anbieter halten 62 % des globalen Cloud-Marktes — europäische Anbieter ~15 %.",
                },
                {
                    "eyebrow": "FINANZ  ·  DORA-FALLE",
                    "title": "Die Bank, die ihre Exit-Strategie nicht kennt",
                    "body": "Seit Januar 2025 verlangt DORA dokumentierte Exit-Strategien für kritische ICT-Drittanbieter. Rechtsgutachten Uni zu Köln (BMI): EU-Rechenzentren reichen nicht — entscheidend ist die rechtliche Beherrschung.",
                },
                {
                    "eyebrow": "HEALTHCARE  ·  PATIENTENDATEN",
                    "title": "Die Onkologin und das Risiko, das niemand sieht",
                    "body": "Eine Münchener Onkologin analysiert mit KI Patientendaten in seltenen Krebserkrankungen. Bei einer CLOUD-Act-Anordnung können die Daten den EU-Rechtsraum verlassen — ohne Information der Klinik, ohne Information der Patientinnen.",
                },
                {
                    "eyebrow": "PUBLIC SECTOR  ·  RECHENSCHAFT",
                    "title": "Der Richter und die Black Box",
                    "body": "Der EU AI Act fordert für Hochrisiko-KI vollständige Rückverfolgbarkeit (Art. 12, 13). Eine deutsche Verwaltung mit proprietären US-Modellen kann gegenüber einem Verwaltungsgericht nicht im Detail erklären, wie eine Entscheidung zustande kam.",
                },
            ],
        ),

        # Akt II
        lambda: build_act_divider(
            prs,
            "II",
            "Die Landkarte.",
            "Souveränität, Compliance-Welle, Säulen, das europäische Anbieter-Ökosystem.",
        ),
        lambda: build_triade(
            prs,
            "AKT II — DIE LANDKARTE  ·  SZENE 1",
            "Souveränität ist nicht Autarkie. Souveränität ist Agency.",
            'Die häufigste Fehlannahme: „Souverän sein heißt, alles selbst zu machen." '
            "Souveränität in einer hypervernetzten Welt heißt etwas anderes — "
            "die Freiheit, sich auf andere zu verlassen, ohne die eigene "
            "Handlungsfähigkeit aufzugeben. Agency in drei Dimensionen:",
            [
                {"num": "01", "title": "Wahl der Partner", "body": "Anbieter werden nach eigenen Kriterien gewählt: Rechtssicherheit, Kosten, Leistung, Werte. Nicht durch Pfadabhängigkeit oder technische Hürden vorgegeben."},
                {"num": "02", "title": "Kontrolle der Regeln", "body": "Die Systeme arbeiten nach den eigenen Standards, Werten und regulatorischen Anforderungen — nicht nach den AGB des Anbieters."},
                {"num": "03", "title": "Fähigkeit zum Wechsel", "body": "Wenn ein Anbieter die Preise verdreifacht, die Bedingungen ändert oder einstellt: Migration ist möglich, ohne das digitale Rückgrat zu zertrümmern."},
            ],
        ),
        lambda: build_editorial(
            prs,
            "AKT II — DIE LANDKARTE  ·  SZENE 2",
            "Die regulatorische Welle 2025–2028.",
            [
                (
                    "Was viele als Compliance-Last empfinden, ist in Wahrheit ein strategischer "
                    "Hebel: Die europäische Regulatorik macht Souveränität zum Wettbewerbsvorteil — "
                    "wenn man sie als Designprinzip nutzt, nicht als Hindernis.",
                    {"size": 13, "italic": True, "color": BODY},
                ),
                (
                    "▸  DSGVO Art. 48 (seit 2018) · 20 Mio. EUR / 4 %",
                    {"size": 11, "color": BODY, "space_before": 12},
                ),
                ("Datenübermittlung an Drittstaatsbehörden nur über Rechtshilfeabkommen — CLOUD Act erfüllt dies nicht.", {"size": 10, "color": BODY_MUTED, "line_spacing": 1.45}),
                ("▸  DORA (Jan 2025) · Persönliche Haftung", {"size": 11, "color": BODY, "space_before": 8}),
                ("Finanzunternehmen müssen Konzentrationsrisiko und Exit-Strategien für ICT-Drittanbieter dokumentieren.", {"size": 10, "color": BODY_MUTED, "line_spacing": 1.45}),
                ("▸  EU Data Act, Kap. VII (Sep 2025) · Diverse", {"size": 11, "color": BODY, "space_before": 8}),
                ("EU-Cloud-Anbieter müssen Drittstaaten-Zugriff aktiv anfechten und prüfen.", {"size": 10, "color": BODY_MUTED, "line_spacing": 1.45}),
                ("▸  NIS2 (D · NIS2UmsuCG, Dez 2025) · 10 Mio. EUR / 2 %", {"size": 11, "color": BODY, "space_before": 8}),
                ("~30.000 Unternehmen mit Pflichten zu Risikomanagement, Meldepflichten, Geschäftsführerhaftung.", {"size": 10, "color": BODY_MUTED, "line_spacing": 1.45}),
                ("▸  EU AI Act · GPAI (Aug 2025) · 35 Mio. EUR / 7 %", {"size": 11, "color": BODY, "space_before": 8}),
                ("General-Purpose-AI-Modelle: Transparenzpflichten in Kraft.", {"size": 10, "color": BODY_MUTED, "line_spacing": 1.45}),
                ("▸  EU AI Act · Hochrisiko (Dez 2027, durch Digital Omnibus verschoben) · 35 Mio. EUR / 7 %", {"size": 11, "color": BODY, "space_before": 8}),
                ("Vollständige Konformitätspflichten (Anhang III) — durch Digital Omnibus vom 7. Mai 2026 von Aug 2026 verschoben.", {"size": 10, "color": BODY_MUTED, "line_spacing": 1.45}),
            ],
        ),
        lambda: build_two_col(
            prs,
            "AKT II — DIE LANDKARTE  ·  SZENE 3",
            "Intelligence Taker oder Intelligence Maker?",
            "Hier liegt die strategische Gabelung. Unternehmen müssen sich entscheiden, welche "
            "Rolle sie im KI-Zeitalter einnehmen wollen.",
            left={
                "eyebrow": "PFAD A  ·  MIETEN",
                "title": "Intelligence Taker",
                "accent": ROST,
                "bullets": [
                    "Konsumiert KI als Standardprodukt über API — schnell, bequem, niedrige Einstiegshürde.",
                    "Die zugrunde liegende Intelligenz bleibt austauschbar. Kein nachhaltiges Differenzierungsmerkmal.",
                    "Eigene Prompts und Daten fließen oft (legitimiert durch AGB) ins Training der nächsten Modell-Generation.",
                    "Die verbesserte Intelligenz kann der Anbieter morgen direkt an den Wettbewerber verkaufen.",
                    "Preis- und Roadmap-Hoheit liegt vollständig beim Anbieter.",
                ],
            },
            right={
                "eyebrow": "PFAD B  ·  BESITZEN",
                "title": "Intelligence Maker",
                "accent": SALBEI,
                "bullets": [
                    "Nutzt offene Modelle (Mistral, Llama, Falcon) auf eigener oder zertifizierter EU-Infrastruktur.",
                    "Veredelt sie durch Fine-Tuning mit eigenen, proprietären Daten.",
                    "Das Ergebnis ist exklusives IP — ein geschütztes Asset, das nicht in fremde Hände gelangt.",
                    "Modellgewichte, Trainingslogik, Audit-Trail: alles im eigenen Einflussbereich.",
                    "Wirtschaftlich überlegen ab dem Skalierungs-Break-Even (typisch 12–18 Monate).",
                ],
            },
        ),
        lambda: build_four_models(
            prs,
            "AKT II — DIE LANDKARTE  ·  SZENE 4",
            "Die Sovereign Operating Matrix.",
            "Souverän zu sein heißt nicht zwangsläufig, ein Rechenzentrum zu kaufen. Vier Modelle, je nach Risikoprofil und Reifestufe.",
            [
                {
                    "name": "Souveränes MaaS",
                    "tagline": "Inferenz-Token bei EU-Anbieter buchen",
                    "bullets": ["Mistral via SAP AI Core", "Aleph Alpha PhariaAI", "100 % EU-Datenhaltung", "Pay-per-use"],
                    "fuer": "KMU, Piloten, Inferenz-Workloads",
                },
                {
                    "name": "Souveränes IaaS / Neoclouds",
                    "tagline": "GPU-Stunden bei spezialisierten EU-Anbietern",
                    "bullets": ["IONOS, OVHcloud, plusserver", "Nebius Europe", "Eigenes Hosting von Open-Weight-Modellen", "BSI-C5- oder Gaia-X-zertifiziert"],
                    "fuer": "Mittelstand mit ML-Kompetenz",
                },
                {
                    "name": "Managed Sovereignty",
                    "tagline": "Dedizierte Server in zertifizierten Rechenzentren",
                    "bullets": ["T-Systems, STACKIT", "OVHcloud Bare Metal", 'Keine „Noisy Neighbors"', "Volle physische Kontrolle"],
                    "fuer": "Regulierte Branchen, sensible Workloads",
                },
                {
                    "name": "Colocation / AI Factory",
                    "tagline": "Eigene Hardware im spezialisierten Rechenzentrum",
                    "bullets": ["Eigene NVIDIA-Cluster", "Direct Liquid Cooling", "InfiniBand-Fabric", "Eigene Foundation-Modelle möglich"],
                    "fuer": "Konzerne, hochsensible Anwendungen",
                },
            ],
        ),
        lambda: build_editorial(
            prs,
            "AKT II — DIE LANDKARTE  ·  SZENE 5",
            "Die 4 Säulen und das europäische Ökosystem.",
            [
                (
                    "Souveränität ist multiplikativ — fehlt eine Säule, sinkt der Gesamtgrad gegen Null. "
                    "Die Behauptung, Europa habe keine relevanten KI-Anbieter, ist 2026 nicht mehr haltbar.",
                    {"size": 12, "italic": True, "color": BODY},
                ),
                ("MODELLE", {"size": 10, "color": PETROL, "bold": True, "font": FONT_MONO, "space_before": 12}),
                (
                    "Mistral AI (Paris) · Frontier-Modelle, Mistral Forge On-Premises, SAP-Partnerschaft seit Nov. 2025, "
                    "Mistral Compute (Schweden, 1,2 Mrd. EUR). Aleph Alpha (Heidelberg) · Spezialist für erklärbare KI in "
                    "regulierten Branchen. Black Forest Labs (D) · Bildgenerierung. LightOn (Paris) · energieeffizientes Training.",
                    {"size": 10.5, "color": BODY},
                ),
                ("INFRASTRUKTUR", {"size": 10, "color": PETROL, "bold": True, "font": FONT_MONO, "space_before": 8}),
                (
                    "STACKIT (Schwarz Digits) · BSI-C5-testiert, eigener AI Cloud Stack. IONOS, T-Systems, OVHcloud, plusserver · "
                    "DSGVO-konform, Gaia-X-Gründungsmitglieder. Delos Cloud · Spezialisiert auf öffentliche Verwaltung.",
                    {"size": 10.5, "color": BODY},
                ),
                ("ÖKOSYSTEM", {"size": 10, "color": PETROL, "bold": True, "font": FONT_MONO, "space_before": 8}),
                (
                    'Gaia-X Trust Framework 3.0 „Danube" seit November 2025 — über 180 Data Spaces in Umsetzung. '
                    "EuroStack-Initiative, IPCEI-CIS, EU Data Act fördern interoperable Cloud-Standards.",
                    {"size": 10.5, "color": BODY},
                ),
            ],
        ),

        # Akt III
        lambda: build_act_divider(
            prs,
            "III",
            "Der Spielzug.",
            "Das CSD-Framework, die 90-Tage-Roadmap und der 12-Fragen-Souveränitäts-Check.",
        ),
        lambda: build_csd_framework(prs),
        lambda: build_roadmap(prs),
        lambda: build_check_questions(prs),
        lambda: build_levels(prs),
        lambda: build_closing(prs),
        lambda: build_sources(prs),
    ]

    slides = [b() for b in builders]
    total = len(slides)

    # Footer mit Seitenzahlen
    for i, slide in enumerate(slides, start=1):
        add_footer(slide, i, total)

    prs.save(output)
    size_kb = output.stat().st_size / 1024
    print(
        f"\n✓ {output.relative_to(repo_root)} geschrieben — "
        f"{total} Slides, {size_kb:.1f} KB"
    )


if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        print(f"\n✗ Fehler: {e}", file=sys.stderr)
        raise
